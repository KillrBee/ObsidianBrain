#!/usr/bin/env python3
"""Regenerate binary test fixtures. Committed outputs live next to this file;
rerun via `make fixtures` when they need to change.

Requires: python-docx, python-pptx, openpyxl (in the dev venv).
The PDF is hand-assembled with computed xref offsets — no PDF lib needed.
"""
from __future__ import annotations

import io
import zlib
from pathlib import Path

HERE = Path(__file__).resolve().parent


def make_pdf(path: Path, text: str) -> None:
    """Minimal valid one-page PDF with a text stream and correct xref table."""
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    buf = io.BytesIO()
    buf.write(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(buf.tell())
        buf.write(f"{i} 0 obj\n".encode() + obj + b"\nendobj\n")
    xref_pos = buf.tell()
    buf.write(f"xref\n0 {len(objects) + 1}\n".encode())
    buf.write(b"0000000000 65535 f \n")
    for off in offsets:
        buf.write(f"{off:010d} 00000 n \n".encode())
    buf.write(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n".encode())
    path.write_bytes(buf.getvalue())


def make_docx(path: Path) -> None:
    import docx
    doc = docx.Document()
    doc.add_heading("Sample DOCX for SecondBrain", level=1)
    doc.add_paragraph("Identity resolution requires dynamic survivorship "
                      "across canonical contacts.")
    doc.add_paragraph("This fixture exercises the docx conversion path.")
    doc.save(path)


def make_pptx(path: Path) -> None:
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX for SecondBrain"
    slide.placeholders[1].text = "Slide body: retrieval substrate, not truth."
    prs.save(path)


def make_xlsx(path: Path) -> None:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Metric", "Value"])
    ws.append(["Sample XLSX for SecondBrain", 42])
    ws.append(["Converted rows", 3])
    wb.save(path)


def main() -> None:
    make_pdf(HERE / "sample.pdf", "Sample PDF for SecondBrain retrieval stack")
    make_docx(HERE / "sample.docx")
    make_pptx(HERE / "sample.pptx")
    make_xlsx(HERE / "sample.xlsx")
    (HERE / "sample.html").write_text(
        "<!doctype html><html><head><title>Sample HTML for SecondBrain</title>"
        "</head><body><h1>Sample HTML for SecondBrain</h1>"
        "<p>Hypertext fixture describing identity resolution evidence.</p>"
        "</body></html>\n")
    # Not a PDF at all — conversion must fail cleanly and land in the manifest.
    (HERE / "corrupt.pdf").write_bytes(b"\x00\x01NOT-A-PDF" + zlib.crc32(b"x").to_bytes(4, "big") * 8)
    print("fixtures written to", HERE)


if __name__ == "__main__":
    main()
